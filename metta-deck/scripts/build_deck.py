"""Runner copy-and-substitute pra deck Metta.

Abordagem (registrada em memory `feedback_pptx_copiar_nao_recriar.md`):
1. Abrir `assets/modelo-geral.pptx` como template
2. Manter TODOS os 25 slides na ordem original
3. Substituir SÓ textos placeholder via find/replace IN-PLACE em runs
4. Preservar 100% das configurações: imagens, layouts, posições, fontes embutidas,
   slide masters, theme colors XML, picture cropping, layer ordering

USO (via Claude Code skill):
    python build_deck.py --briefing briefing.json --out proposta-cliente.pptx

USO direto:
    python build_deck.py \\
        --client "ACME" \\
        --statement "Como vamos juntos · ACME." \\
        --cta "AGENDAR COM JOÃO" \\
        --date-month "Junho/2026" \\
        --out proposta-acme.pptx

Briefing JSON (formato completo):
    {
      "client": "LEROY MERLIN",
      "statement": "Como vamos juntos · Leroy Merlin.",
      "cta": "AGENDAR COM CAIQUE",
      "date_month": "Junho/2026",
      "kpi_subtitle": "Os indicadores que sustentam a operação · exemplo prático.",
      "extra_replacements": {
        "Carlos Oliveira": "[Nome real depoimento 1]"
      }
    }
"""
from __future__ import annotations
import argparse
import json
import os
import sys
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("ERRO: instale python-pptx → pip install python-pptx", file=sys.stderr)
    sys.exit(1)


# Resolve template path via env var ${CLAUDE_PLUGIN_ROOT} (plugin context)
# ou fallback pro path relativo (dev local)
def get_plugin_root() -> Path:
    env = os.environ.get("CLAUDE_PLUGIN_ROOT")
    if env:
        return Path(env)
    # Fallback: assume script está em metta-deck/scripts/build_deck.py
    return Path(__file__).resolve().parent.parent


PLUGIN_ROOT = get_plugin_root()
TEMPLATE_PATH = PLUGIN_ROOT / "assets" / "modelo-geral.pptx"


# Placeholders default do PPT Modelo Geral (catálogo validado em 2026-05-25)
DEFAULT_PLACEHOLDERS = {
    "[NOME CLIENTE]":                         None,  # client_upper (obrigatório)
    "Lorem ipsum dolor sit amet":             None,  # statement
    "Grupo Linhares":                         None,  # client_title
    "AGENDAR DIAGNÓSTICO":                    None,  # cta
    "Próximas 2 semanas":                     None,  # date_month
    "Exemplificar na situação do cliente.":   None,  # kpi_subtitle
}


def build_replacements(briefing: dict) -> dict:
    """Constrói dict de substituições a partir do briefing.

    Aplica defaults sensatos quando o briefing não cobre um placeholder.
    """
    client = briefing.get("client", "").strip()
    if not client:
        raise ValueError("briefing.client é obrigatório")

    replacements = {
        "[NOME CLIENTE]":    client.upper(),
        "Grupo Linhares":    client.title(),
        "AGENDAR DIAGNÓSTICO": briefing.get("cta", "AGENDAR DIAGNÓSTICO").upper(),
        "Próximas 2 semanas":  briefing.get("date_month", "Próximas 2 semanas"),
        "Lorem ipsum dolor sit amet": briefing.get(
            "statement",
            f"Como vamos juntos · {client.title()}."
        ),
        "Exemplificar na situação do cliente.": briefing.get(
            "kpi_subtitle",
            "Os indicadores que sustentam a operação comercial · "
            "exemplo prático do método."
        ),
    }

    # Extras: substituições específicas do briefing (ex: nomes depoimentos)
    extras = briefing.get("extra_replacements") or {}
    replacements.update(extras)

    return replacements


def replace_in_run(run, old: str, new: str) -> bool:
    """Substitui SE o texto do run contém `old`. Preserva formatação do run."""
    if old in run.text:
        run.text = run.text.replace(old, new)
        return True
    return False


def replace_in_slide(slide, replacements: dict) -> dict:
    """Aplica substituições em todos os runs do slide."""
    stats = {}
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                for old, new in replacements.items():
                    if replace_in_run(run, old, new):
                        stats[old] = stats.get(old, 0) + 1
    return stats


def build(briefing: dict, out_path: Path, verbose: bool = True) -> Path:
    """Gera o PPTX a partir do template + briefing."""
    if not TEMPLATE_PATH.exists():
        raise FileNotFoundError(f"Template não encontrado: {TEMPLATE_PATH}")

    if verbose:
        print(f"Carregando template: {TEMPLATE_PATH.name}")
    prs = Presentation(str(TEMPLATE_PATH))
    if verbose:
        print(f"  → {len(prs.slides)} slides · "
              f"{prs.slide_width/914400:.2f}x{prs.slide_height/914400:.2f}in")

    replacements = build_replacements(briefing)

    total = {}
    for idx, slide in enumerate(prs.slides):
        stats = replace_in_slide(slide, replacements)
        if stats and verbose:
            summary = ", ".join(f"{k[:30]}×{v}" for k, v in stats.items())
            print(f"  slide {idx+1:2d}: {summary}")
        for k, v in stats.items():
            total[k] = total.get(k, 0) + v

    if verbose:
        print("\nResumo:")
        for k, v in total.items():
            print(f"  '{k[:45]}'  × {v}")

    out_path = Path(out_path).resolve()
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    if verbose:
        print(f"\nOK · saved: {out_path}")
    return out_path


def main():
    ap = argparse.ArgumentParser(description="Gerador de deck Metta")
    ap.add_argument("--briefing", type=str,
                     help="Path pra arquivo JSON com briefing completo")
    ap.add_argument("--client", type=str, help="Nome do cliente (ex: LEROY MERLIN)")
    ap.add_argument("--statement", type=str, help="Statement da capa")
    ap.add_argument("--cta", type=str, help="Texto do CTA (ex: AGENDAR COM JOÃO)")
    ap.add_argument("--date-month", type=str,
                     help="Próximos passos (ex: Junho/2026)")
    ap.add_argument("--kpi-subtitle", type=str, help="Subtitle slide 11 KPIs")
    ap.add_argument("--out", type=str, required=True,
                     help="Path do .pptx de saída")
    ap.add_argument("--quiet", action="store_true", help="Suprime logs")
    args = ap.parse_args()

    if args.briefing:
        briefing = json.loads(Path(args.briefing).read_text(encoding="utf-8"))
    else:
        if not args.client:
            ap.error("--briefing OU --client (com demais flags) é obrigatório")
        briefing = {
            "client": args.client,
            "statement": args.statement,
            "cta": args.cta,
            "date_month": args.date_month,
            "kpi_subtitle": args.kpi_subtitle,
        }
        briefing = {k: v for k, v in briefing.items() if v is not None}

    build(briefing, Path(args.out), verbose=not args.quiet)


if __name__ == "__main__":
    main()
