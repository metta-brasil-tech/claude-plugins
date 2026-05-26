"""Lib canônica de builders Metta — fonte única pra construir slides PPTX.

Importe daqui SEMPRE. Não duplique helpers em scripts dispersos.

Convenção:
- new_deck() → cria Presentation 13.33×7.5 widescreen
- save(prs, path) → salva
- add_header / add_footer → boilerplate canônico
- add_title_split → title com primeira letra bold (truque editorial Metta)
- slide_* → cada tipo de slide canônico

Fonte de verdade: design/metta-tokens.md + arquitetura/prompt-pptx-metta-master.md
Modelado em: PPT Modelo Geral.pptx (25 slides referência)
"""
from __future__ import annotations

import math
from pathlib import Path
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION, XL_TICK_MARK
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn


# ============================================================
# TOKENS — design system Metta v3.0 (FIEL ao PPT Modelo Geral)
# Fonte: design/metta-pptx-canonico.md §2
# ============================================================
# Backgrounds
BG_DARK       = RGBColor(0x0C, 0x16, 0x1B)   # bg-dark · slide escuro padrão
BG_LIGHT      = RGBColor(0xFF, 0xFF, 0xFF)   # bg-light · branco puro padrão
BG_LIGHT_COOL = RGBColor(0xF2, 0xF7, 0xFC)   # bg-light-cool · off-white frio (capa light)

# Containers (cards)
CONTAINER_DARK  = RGBColor(0x13, 0x1F, 0x25)  # cards sobre dark
CONTAINER_LIGHT = RGBColor(0xEB, 0xF3, 0xF7)  # cards sobre light (#EBF3F7)

# Texto sobre dark
TEXT_DARK_BOLD   = RGBColor(0xF2, 0xF7, 0xFC)  # off-white pra body bold em dark
TEXT_DARK_PURE   = RGBColor(0xFF, 0xFF, 0xFF)  # branco puro · uso restrito (big number unit, divider title)
TEXT_DARK_MUTED  = RGBColor(0xB0, 0xCA, 0xD8)  # body/eyebrow/footer em dark

# Texto sobre light
TEXT_LIGHT_BOLD  = RGBColor(0x0C, 0x16, 0x1B)  # bold/title em light
TEXT_LIGHT_MUTED = RGBColor(0x43, 0x59, 0x65)  # body/eyebrow/footer em light
TEXT_LIGHT_SOFT  = RGBColor(0x11, 0x1E, 0x25)  # variante leve (subline capa light)

# Yellow
YELLOW       = RGBColor(0xFF, 0xBE, 0x18)   # ★ brand primary
YELLOW_SOFT  = RGBColor(0xFF, 0xE4, 0xA1)   # uso RARO (antes/depois)

# Charts data
GREEN_POS    = RGBColor(0x2E, 0x7D, 0x32)
RED_NEG      = RGBColor(0xC6, 0x28, 0x28)

# ALIASES retrocompat (v2 → v3) — não usar em código novo
NIGHT_5      = BG_DARK
NIGHT_10     = TEXT_LIGHT_BOLD     # = BG_DARK no hex (0C161B)
NIGHT_15     = CONTAINER_DARK
NIGHT_40     = TEXT_LIGHT_MUTED
NIGHT_70     = RGBColor(0x94, 0xB5, 0xC8)
NIGHT_85     = TEXT_DARK_MUTED
NIGHT_95     = CONTAINER_LIGHT
NIGHT_100    = BG_LIGHT
YELLOW_WHISP = RGBColor(0xFF, 0xFA, 0xEC)

# Tipografia (Zalando+Inter — sem SF Pro em PPTX)
FONT_HEADS         = "Zalando Sans Expanded"
FONT_HEADS_SEMI    = "Zalando Sans SemiExpanded"  # raro · subhead nível 2
FONT_BODY          = "Inter"

# Canvas 16:9 widescreen
W_IN = 13.333
H_IN = 7.5
W = Inches(W_IN)
H = Inches(H_IN)

ASSETS = Path(__file__).parent / "assets"
LOGO_DARK_PATH  = ASSETS / "logo_dark.png"
LOGO_LIGHT_PATH = ASSETS / "logo_light.png"


# ============================================================
# DECK MGMT
# ============================================================
def new_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs


def save(prs: Presentation, path: str | Path):
    prs.save(str(path))


def new_slide(prs, *, theme="light", bg=None):
    """Cria slide em branco com background do tema.

    theme: 'light' (BG_LIGHT) | 'dark' (BG_DARK) | 'light-cool' (BG_LIGHT_COOL) | 'yellow'
    bg: override explícito (RGBColor) — vence sobre theme
    """
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    if bg is None:
        bg = {
            "light": BG_LIGHT,
            "dark": BG_DARK,
            "light-cool": BG_LIGHT_COOL,
            "yellow": YELLOW,
        }.get(theme, BG_LIGHT)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


# ============================================================
# HELPERS — anti-shadow + primitives
# ============================================================
def _strip_shadow(shape):
    """Remove shadow default do shape via XML."""
    try:
        sppr = shape._element.find(qn("p:spPr"))
        if sppr is None:
            return
        for effect in sppr.findall(qn("a:effectLst")):
            sppr.remove(effect)
        etree.SubElement(sppr, qn("a:effectLst"))
    except Exception:
        pass


def _strip_chart_shadows(chart):
    """Remove shadow + 3D effects do chart."""
    for el in chart._chartSpace.iter():
        tag = etree.QName(el).localname
        if tag in ("plotArea", "ser", "categoryAxis", "valAxis",
                   "lineChart", "barChart", "doughnutChart", "areaChart"):
            for sub in el.findall(qn("c:spPr")):
                for old in sub.findall(qn("a:effectLst")):
                    sub.remove(old)
                etree.SubElement(sub, qn("a:effectLst"))


def add_rect(slide, x, y, w, h, fill, line=None, radius=0):
    """Retângulo (opcionalmente arredondado), sem shadow."""
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius > 0 else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shp_type, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    if radius > 0:
        try:
            shape.adjustments[0] = radius
        except Exception:
            pass
    _strip_shadow(shape)
    return shape


def add_ellipse(slide, x, y, w, h, fill=None, line=None, line_w=1):
    """Elipse/círculo sem shadow."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    _strip_shadow(shape)
    return shape


def add_line(slide, x1, y1, x2, y2, color, weight=0.75):
    """Linha simples sem shadow."""
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    _strip_shadow(line)
    return line


def add_text(slide, x, y, w, h, text, *,
             font_name=None, size=14, bold=False, color=NIGHT_10,
             align="left", anchor="top", line_spacing=1.2, letter_spacing=0,
             is_head=False):
    """Caixa de texto formatada. is_head=True usa Zalando, senão Inter."""
    if font_name is None:
        font_name = FONT_HEADS if is_head else FONT_BODY
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
                           "bottom": MSO_ANCHOR.BOTTOM}[anchor]
    for i, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                       "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        if letter_spacing:
            run._r.get_or_add_rPr().set("spc", str(letter_spacing))
    _strip_shadow(tb)
    return tb


def add_title_split(slide, *, x, y, w, text, size=30, color=NIGHT_10):
    """Title com PRIMEIRA LETRA bold + resto regular (truque editorial Metta).

    Ex: "Indicadores que olhamos" → 'I' (bold) + 'ndicadores que olhamos' (regular)
    """
    tb = slide.shapes.add_textbox(x, y, w, Inches(0.5))
    tf = tb.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = 1.0

    # v3: parsing markdown-like
    #   - sem `**word**` → tudo bold (split apenas estrutural primeira letra)
    #   - com `**word**` → palavras-chave bold + resto regular (truque editorial)
    text = str(text)
    if "**" in text:
        # Parse: split por ** mantendo flags bold/regular
        parts = text.split("**")
        # parts alternam: [regular, bold, regular, bold, ...]
        for i, part in enumerate(parts):
            if not part:
                continue
            is_bold = (i % 2 == 1)  # ímpares são bold
            r = p.add_run()
            r.text = part
            r.font.name = FONT_HEADS
            r.font.size = Pt(size)
            r.font.bold = is_bold
            r.font.color.rgb = color
    else:
        # Default: tudo bold, primeira letra em run separado (controle de tracking)
        if len(text) >= 1:
            r1 = p.add_run()
            r1.text = text[0]
            r1.font.name = FONT_HEADS
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.color.rgb = color
            if len(text) > 1:
                r2 = p.add_run()
                r2.text = text[1:]
                r2.font.name = FONT_HEADS
                r2.font.size = Pt(size)
                r2.font.bold = True
                r2.font.color.rgb = color
    _strip_shadow(tb)
    return tb


def add_logo(slide, theme="light", *, mode="main"):
    """Logo Metta no canto superior esquerdo.

    mode:
      - 'main' (default): x=0.4 y=0.317 w=1.067 h=0.32 — slides padrão
      - 'mini': x=0.556 y=0.389 w=0.903 h=0.25 — capa, splits, divider editorial
      - 'footer-mini': x=0.556 y=7.028 w=0.653 h=0.181 — bottom em slides foto bleed
    """
    path = LOGO_DARK_PATH if theme == "light" else LOGO_LIGHT_PATH
    if not path.exists():
        return None
    if mode == "mini":
        return slide.shapes.add_picture(str(path), Inches(0.556), Inches(0.389),
                                         width=Inches(0.903), height=Inches(0.25))
    elif mode == "footer-mini":
        return slide.shapes.add_picture(str(path), Inches(0.556), Inches(7.028),
                                         width=Inches(0.653), height=Inches(0.181))
    # main (default)
    return slide.shapes.add_picture(str(path), Inches(0.4), Inches(0.317),
                                     width=Inches(1.067), height=Inches(0.32))


# ============================================================
# HEADER + FOOTER (canônicos — usar em todo slide de conteúdo)
# ============================================================
def add_header(slide, eyebrow="", client="[NOME CLIENTE]", *,
                theme="light", mode="main"):
    """Header canônico Metta — 3 modos (modelo geral §4 canônico).

    mode:
      - 'main': slide padrão de conteúdo (logo + eyebrow + cliente + divider full)
      - 'compact': capa / split (logo mini + date direita 6pt; sem divider)
      - 'micro': foto bleed editorial (só eyebrow micro 8pt sobre foto)
      - 'split': quem somos com coluna estreita (logo + eyebrow + cliente + divider PARCIAL w=4.385)
    """
    muted = TEXT_LIGHT_MUTED if theme == "light" else TEXT_DARK_MUTED
    divider_color = NIGHT_85 if theme == "light" else NIGHT_40

    if mode == "main":
        add_logo(slide, theme=theme, mode="main")
        add_text(slide, Inches(2.0), Inches(0.32), Inches(7.0), Inches(0.25),
                 text=eyebrow.upper(), size=9, bold=True, color=muted,
                 is_head=True, anchor="middle")
        add_text(slide, Inches(8.0), Inches(0.32), Inches(4.9), Inches(0.25),
                 text=client.upper(), size=9, bold=True, color=muted,
                 align="right", is_head=True, anchor="middle")
        add_line(slide, Inches(0.4), Inches(0.72), Inches(12.9), Inches(0.72),
                 divider_color, weight=0.5)

    elif mode == "compact":
        # Capa: só logo mini + data micro top-right
        add_logo(slide, theme=theme, mode="mini")
        add_text(slide, Inches(11.944), Inches(0.444), Inches(0.584), Inches(0.111),
                 text="MAIO 2026", size=6, bold=True, color=muted,
                 is_head=True, anchor="top")

    elif mode == "micro":
        # Foto bleed editorial: só eyebrow micro
        text_color = TEXT_DARK_PURE if theme == "dark" else TEXT_DARK_PURE  # branco sobre foto
        add_text(slide, Inches(0.556), Inches(0.389), Inches(1.646), Inches(0.132),
                 text=eyebrow.upper(), size=8, bold=True, color=text_color,
                 is_head=True, anchor="top")

    elif mode == "split":
        # Quem somos coluna estreita
        add_logo(slide, theme=theme, mode="main")  # logo na esquerda
        # Eyebrow ocupa o lado esquerdo apenas (até 4.385)
        add_text(slide, Inches(2.045), Inches(0.362), Inches(2.0), Inches(0.152),
                 text=eyebrow.upper(), size=9, bold=True, color=muted,
                 is_head=True, anchor="middle")
        # Divider parcial só esquerdo
        add_line(slide, Inches(0.4), Inches(0.72), Inches(4.785), Inches(0.72),
                 divider_color, weight=0.5)


def add_footer(slide, *, date="MAIO 2026", client="[NOME CLIENTE]", section="",
                theme="light", mode="main"):
    """Footer canônico Metta — 3 modos (modelo geral §5 canônico).

    mode:
      - 'main': h=0.3 size=8pt — slide padrão
      - 'compact': h=0.135 size=8pt — capa/split
      - 'micro': h=0.111 size=6pt bold — slide foto editorial bleed
      - 'split': h=0.135 só esquerda (sem section/arrow) — quem somos split column
    """
    muted = TEXT_LIGHT_MUTED if theme == "light" else TEXT_DARK_MUTED
    divider_color = NIGHT_85 if theme == "light" else NIGHT_40
    arrow_color = YELLOW if theme == "light" else TEXT_DARK_BOLD

    if mode == "main":
        add_line(slide, Inches(0.4), Inches(7.0), Inches(12.9), Inches(7.0),
                 divider_color, weight=0.5)
        add_text(slide, Inches(0.4), Inches(7.15), Inches(3.0), Inches(0.3),
                 text=date, size=8, color=muted, is_head=True, anchor="top")
        add_text(slide, Inches(5.5), Inches(7.15), Inches(3.0), Inches(0.3),
                 text=client, size=8, color=muted, align="center",
                 is_head=True, anchor="top")
        add_text(slide, Inches(10.0), Inches(7.15), Inches(2.5), Inches(0.3),
                 text=section.upper(), size=8, color=muted, align="right",
                 is_head=True, anchor="top")
        add_text(slide, Inches(12.6), Inches(7.15), Inches(0.3), Inches(0.3),
                 text="↗", size=10, bold=True, color=arrow_color, align="right",
                 is_head=True)

    elif mode == "compact":
        add_line(slide, Inches(0.4), Inches(7.0), Inches(12.9), Inches(7.0),
                 divider_color, weight=0.5)
        add_text(slide, Inches(0.4), Inches(7.15), Inches(3.0), Inches(0.135),
                 text=date, size=8, color=muted, is_head=True, anchor="top")
        add_text(slide, Inches(5.5), Inches(7.15), Inches(3.0), Inches(0.135),
                 text=client, size=8, color=muted, align="center",
                 is_head=True, anchor="top")
        add_text(slide, Inches(10.0), Inches(7.15), Inches(2.5), Inches(0.135),
                 text=section.upper(), size=8, color=muted, align="right",
                 is_head=True, anchor="top")
        add_text(slide, Inches(12.6), Inches(7.15), Inches(0.3), Inches(0.168),
                 text="↗", size=10, bold=True, color=arrow_color, align="right",
                 is_head=True)

    elif mode == "micro":
        # Foto bleed: footer micro 6pt, logo mini bottom-left, arrow → (não ↗)
        add_logo(slide, theme=theme, mode="footer-mini")
        white = TEXT_DARK_PURE
        add_text(slide, Inches(3.75), Inches(7.069), Inches(0.584), Inches(0.111),
                 text=date.title(), size=6, bold=True, color=white,
                 is_head=True, anchor="top")
        add_text(slide, Inches(6.528), Inches(7.069), Inches(0.812), Inches(0.111),
                 text=client.title(), size=6, bold=True, color=TEXT_LIGHT_MUTED,
                 is_head=True, anchor="top")
        add_text(slide, Inches(9.305), Inches(7.069), Inches(0.639), Inches(0.111),
                 text=section.title(), size=6, bold=True, color=TEXT_LIGHT_MUTED,
                 is_head=True, anchor="top")
        add_text(slide, Inches(12.514), Inches(7.007), Inches(0.132), Inches(0.188),
                 text="→", size=10, bold=True, color=white, is_head=True)

    elif mode == "split":
        # Quem somos coluna estreita: divider parcial + só date + cliente
        add_line(slide, Inches(0.4), Inches(7.0), Inches(4.589), Inches(7.0),
                 divider_color, weight=0.5)
        add_text(slide, Inches(0.4), Inches(7.15), Inches(3.0), Inches(0.135),
                 text=date, size=8, color=muted, is_head=True, anchor="top")
        add_text(slide, Inches(1.583), Inches(7.15), Inches(2.5), Inches(0.135),
                 text=section.upper() or client.upper(), size=8, color=muted,
                 align="right", is_head=True, anchor="top")


def add_title_block(slide, *, title, theme="light", y=1.10, h=0.404, size=30, align="left"):
    """Title canônico Metta v3 (FIEL ao modelo geral §3 canônico):
    - y=1.10, h=0.404 (não 0.5 ou 0.7!)
    - size 30pt default (slide content padrão)
    - sem linha yellow accent abaixo
    - parsing markdown `**bold**` no title via add_title_split
    """
    fg = TEXT_LIGHT_BOLD if theme == "light" else TEXT_DARK_BOLD
    add_title_split(slide, x=Inches(0.4), y=Inches(y),
                     w=Inches(12.5), text=title, size=size, color=fg)


def standard_frame(slide, *, eyebrow, title, section, client="[NOME CLIENTE]",
                    date="MAIO 2026", theme="light"):
    """Header MAIN + title + footer MAIN. Boilerplate slide padrão de conteúdo."""
    add_header(slide, eyebrow=eyebrow, client=client, theme=theme, mode="main")
    add_title_block(slide, title=title, theme=theme)
    add_footer(slide, date=date, client=client, section=section, theme=theme, mode="main")


# ============================================================
# CHART HELPERS
# ============================================================
def _style_chart_axes(chart, *, theme="light"):
    """McKinsey-style: sem axis lines, gridlines sutis, ticks ocultos."""
    muted = NIGHT_40 if theme == "light" else NIGHT_85
    gridline_hex = "B0CAD8" if theme == "light" else "435965"

    for axis in [chart.category_axis, chart.value_axis]:
        try:
            axis.format.line.fill.background()
            axis.tick_labels.font.name = FONT_BODY
            axis.tick_labels.font.size = Pt(9)
            axis.tick_labels.font.color.rgb = muted
            axis.major_tick_mark = XL_TICK_MARK.NONE
            axis.minor_tick_mark = XL_TICK_MARK.NONE
        except Exception:
            pass

    try:
        gridlines_el = chart.value_axis._element.find(qn("c:majorGridlines"))
        if gridlines_el is not None:
            for old in gridlines_el.findall(qn("c:spPr")):
                gridlines_el.remove(old)
            sppr = etree.SubElement(gridlines_el, qn("c:spPr"))
            ln = etree.SubElement(sppr, qn("a:ln"))
            ln.set("w", "6350")
            solid = etree.SubElement(ln, qn("a:solidFill"))
            srgb = etree.SubElement(solid, qn("a:srgbClr"))
            srgb.set("val", gridline_hex)
            alpha = etree.SubElement(srgb, qn("a:alpha"))
            alpha.set("val", "50000")
    except Exception:
        pass

    _strip_chart_shadows(chart)


# ============================================================
# NOVOS BUILDERS — Frameworks Visuais (não existiam no PPT Modelo)
# ============================================================
def slide_matriz_2x2(prs, *, title, x_label, y_label,
                       x_low_label, x_high_label, y_low_label, y_high_label,
                       items, winner_quadrant="top-right",
                       client="[NOME CLIENTE]", date="MAIO 2026", theme="light"):
    """Matriz 2×2 de priorização.

    EXPERIMENTAL — v1 não passou na validação visual do user (removido do deck canônico
    em 2026-05-25). Bubbles podem ficar cobertas pelo border do quadrant winner.
    Não usar em deck cliente sem revisão manual.

    items: lista de {"name": str, "x": 0-1, "y": 0-1, "highlight": bool}
    winner_quadrant: 'top-right' | 'top-left' | 'bottom-right' | 'bottom-left'
    """
    slide = new_slide(prs, theme=theme)
    standard_frame(slide, eyebrow="MATRIZ · PRIORIZAÇÃO", title=title,
                    section="Estratégia", client=client, date=date, theme=theme)

    fg = NIGHT_10 if theme == "light" else NIGHT_100
    container_bg = NIGHT_95 if theme == "light" else NIGHT_15
    muted = NIGHT_40 if theme == "light" else NIGHT_85

    # Área da matriz
    mx, my, mw, mh = 2.5, 2.3, 8.5, 4.3  # x, y, width, height (inches)

    # Quadrante highlight (winner) — background yellow soft
    quad_w = mw / 2
    quad_h = mh / 2
    quad_pos = {
        "top-right":     (mx + quad_w, my),
        "top-left":      (mx, my),
        "bottom-right":  (mx + quad_w, my + quad_h),
        "bottom-left":   (mx, my + quad_h),
    }
    qx, qy = quad_pos[winner_quadrant]
    add_rect(slide, Inches(qx), Inches(qy), Inches(quad_w), Inches(quad_h),
             fill=YELLOW_SOFT, radius=0)

    # Linhas do eixo
    cross_x = mx + mw / 2
    cross_y = my + mh / 2
    # Linha vertical center
    add_line(slide, Inches(cross_x), Inches(my), Inches(cross_x), Inches(my + mh),
             NIGHT_85 if theme == "light" else NIGHT_40, weight=1)
    # Linha horizontal center
    add_line(slide, Inches(mx), Inches(cross_y), Inches(mx + mw), Inches(cross_y),
             NIGHT_85 if theme == "light" else NIGHT_40, weight=1)
    # Borda externa
    add_rect(slide, Inches(mx), Inches(my), Inches(mw), Inches(mh),
             fill=NIGHT_100 if theme == "light" else NIGHT_10,
             line=NIGHT_85 if theme == "light" else NIGHT_40, radius=0)
    # IMPORTANTE: rect tem fill — vai cobrir o quadrante yellow. Refazer ordem:
    # já fizemos quad antes, agora apenas border externo sem fill — refazer:
    # (python-pptx não tem stack order trivial, então faremos via XML adjustment)

    # Labels axis
    # Y axis label (vertical à esquerda)
    add_text(slide, Inches(0.4), Inches(my + mh / 2 - 0.2), Inches(2.0), Inches(0.4),
             text=y_label.upper(), size=11, bold=True, color=muted,
             letter_spacing=300, is_head=True, align="center")
    # Y axis low/high labels
    add_text(slide, Inches(mx - 1.1), Inches(my - 0.1), Inches(1.0), Inches(0.3),
             text=y_high_label.upper(), size=8, bold=True, color=muted,
             letter_spacing=200, is_head=True, align="right")
    add_text(slide, Inches(mx - 1.1), Inches(my + mh - 0.2), Inches(1.0), Inches(0.3),
             text=y_low_label.upper(), size=8, bold=True, color=muted,
             letter_spacing=200, is_head=True, align="right")

    # X axis label (centro embaixo)
    add_text(slide, Inches(mx), Inches(my + mh + 0.3), Inches(mw), Inches(0.4),
             text=x_label.upper(), size=11, bold=True, color=muted,
             letter_spacing=300, is_head=True, align="center")
    add_text(slide, Inches(mx - 0.5), Inches(my + mh + 0.05), Inches(1.5), Inches(0.3),
             text=x_low_label.upper(), size=8, bold=True, color=muted,
             letter_spacing=200, is_head=True, align="left")
    add_text(slide, Inches(mx + mw - 1.0), Inches(my + mh + 0.05), Inches(1.5), Inches(0.3),
             text=x_high_label.upper(), size=8, bold=True, color=muted,
             letter_spacing=200, is_head=True, align="right")

    # Items (bubbles posicionadas)
    for item in items:
        # Posição absoluta na matriz
        ix = mx + item["x"] * mw
        iy = my + (1 - item["y"]) * mh  # invertido pq Y do screen é top-down
        size = item.get("size", 0.7)  # diameter em inches
        highlight = item.get("highlight", False)
        fill = YELLOW if highlight else (NIGHT_85 if theme == "light" else NIGHT_40)
        text_color = NIGHT_10 if highlight else (NIGHT_100 if theme == "dark" else NIGHT_100)

        add_ellipse(slide, Inches(ix - size / 2), Inches(iy - size / 2),
                     Inches(size), Inches(size), fill=fill)
        add_text(slide, Inches(ix - 1.0), Inches(iy + size / 2 + 0.05),
                 Inches(2.0), Inches(0.4),
                 text=item["name"], size=10, bold=highlight, color=fg,
                 align="center", line_spacing=1.2, is_head=True)


def slide_hub_spoke(prs, *, title, center_label, satellites,
                      client="[NOME CLIENTE]", date="MAIO 2026", theme="light"):
    """Hub central + 4-6 satellites conectados.

    EXPERIMENTAL — v1 não passou na validação visual do user (removido do deck canônico
    em 2026-05-25). Linhas retas + tipografia podem ficar confusas com 6+ satellites.
    Não usar em deck cliente sem revisão manual.

    satellites: lista de {"label": str, "body": str (opcional)}
    """
    slide = new_slide(prs, theme=theme)
    standard_frame(slide, eyebrow="ECOSSISTEMA · ESTRUTURA", title=title,
                    section="Método", client=client, date=date, theme=theme)

    fg = NIGHT_10 if theme == "light" else NIGHT_100
    container_bg = NIGHT_95 if theme == "light" else NIGHT_15
    muted = NIGHT_40 if theme == "light" else NIGHT_85
    line_color = NIGHT_85 if theme == "light" else NIGHT_40

    # Centro
    cx, cy = 6.667, 4.35  # canvas center
    hub_r = 1.0
    # Hub circle yellow
    add_ellipse(slide, Inches(cx - hub_r), Inches(cy - hub_r),
                 Inches(hub_r * 2), Inches(hub_r * 2), fill=YELLOW)
    add_text(slide, Inches(cx - hub_r), Inches(cy - 0.4),
             Inches(hub_r * 2), Inches(0.8),
             text=center_label, size=14, bold=True, color=NIGHT_10,
             align="center", anchor="middle", is_head=True, line_spacing=1.15)

    # Satellites — círculos ao redor em ângulos uniformes
    n = len(satellites)
    radius = 2.6  # distance from center
    sat_r = 0.65
    start_angle = -math.pi / 2  # começa no topo (12 horas)
    angles = [start_angle + 2 * math.pi * i / n for i in range(n)]

    for i, (sat, angle) in enumerate(zip(satellites, angles)):
        sx = cx + radius * math.cos(angle)
        sy = cy + radius * math.sin(angle)

        # Linha conectora (centro → satélite)
        # Calcula ponto de borda do hub e do satélite
        hub_edge_x = cx + hub_r * math.cos(angle)
        hub_edge_y = cy + hub_r * math.sin(angle)
        sat_edge_x = sx - sat_r * math.cos(angle)
        sat_edge_y = sy - sat_r * math.sin(angle)
        add_line(slide, Inches(hub_edge_x), Inches(hub_edge_y),
                 Inches(sat_edge_x), Inches(sat_edge_y), line_color, weight=1.2)

        # Satellite circle
        add_ellipse(slide, Inches(sx - sat_r), Inches(sy - sat_r),
                     Inches(sat_r * 2), Inches(sat_r * 2),
                     fill=container_bg, line=line_color, line_w=1)
        # Label dentro
        add_text(slide, Inches(sx - sat_r), Inches(sy - 0.25),
                 Inches(sat_r * 2), Inches(0.5),
                 text=sat["label"], size=10, bold=True, color=fg,
                 align="center", anchor="middle", is_head=True, line_spacing=1.1)


def slide_timeline(prs, *, title, events, client="[NOME CLIENTE]",
                    date="MAIO 2026", theme="light"):
    """Timeline horizontal com 4-8 marcos.

    events: lista de {"date": str, "title": str, "body": str (opcional), "highlight": bool}
    """
    slide = new_slide(prs, theme=theme)
    standard_frame(slide, eyebrow="LINHA DO TEMPO · MARCOS", title=title,
                    section="Roadmap", client=client, date=date, theme=theme)

    fg = NIGHT_10 if theme == "light" else NIGHT_100
    muted = NIGHT_40 if theme == "light" else NIGHT_85
    container_bg = NIGHT_95 if theme == "light" else NIGHT_15

    # Timeline line
    n = len(events)
    timeline_y = 4.0
    start_x = 1.0
    end_x = 12.333
    width = end_x - start_x
    add_line(slide, Inches(start_x), Inches(timeline_y),
             Inches(end_x), Inches(timeline_y),
             NIGHT_85 if theme == "light" else NIGHT_40, weight=2)

    step = width / (n - 1) if n > 1 else 0
    for i, ev in enumerate(events):
        ex = start_x + i * step
        highlight = ev.get("highlight", False)

        # Dot no eixo
        dot_r = 0.18 if highlight else 0.13
        dot_color = YELLOW if highlight else (NIGHT_40 if theme == "light" else NIGHT_85)
        add_ellipse(slide, Inches(ex - dot_r), Inches(timeline_y - dot_r),
                     Inches(dot_r * 2), Inches(dot_r * 2), fill=dot_color)

        # Alternância: ímpares ACIMA, pares ABAIXO
        # v2: hierarquia tight — DATE sempre adjacente ao eixo, TITLE+BODY se afastando
        above = i % 2 == 0
        if above:
            # ACIMA (lê do eixo pra fora): TITLE topo > BODY > DATE adjacente
            add_text(slide, Inches(ex - 1.2), Inches(timeline_y - 1.29),
                     Inches(2.4), Inches(0.22),
                     text=ev["title"], size=13, bold=True, color=fg,
                     align="center", is_head=True, line_spacing=1.15)
            if ev.get("body"):
                add_text(slide, Inches(ex - 1.2), Inches(timeline_y - 1.00),
                         Inches(2.4), Inches(0.15),
                         text=ev["body"], size=9, color=muted, align="center",
                         line_spacing=1.3)
            add_text(slide, Inches(ex - 1.0), Inches(timeline_y - 0.72),
                     Inches(2.0), Inches(0.17),
                     text=ev["date"].upper(), size=10, bold=True, color=YELLOW,
                     align="center", letter_spacing=200, is_head=True)
            add_line(slide, Inches(ex), Inches(timeline_y - dot_r),
                     Inches(ex), Inches(timeline_y - 0.49),
                     muted, weight=0.5)
        else:
            # ABAIXO (lê do eixo pra fora): DATE adjacente > TITLE > BODY
            add_text(slide, Inches(ex - 1.0), Inches(timeline_y + 0.52),
                     Inches(2.0), Inches(0.17),
                     text=ev["date"].upper(), size=10, bold=True, color=YELLOW,
                     align="center", letter_spacing=200, is_head=True)
            add_text(slide, Inches(ex - 1.2), Inches(timeline_y + 0.74),
                     Inches(2.4), Inches(0.22),
                     text=ev["title"], size=13, bold=True, color=fg,
                     align="center", is_head=True, line_spacing=1.15)
            if ev.get("body"):
                add_text(slide, Inches(ex - 1.2), Inches(timeline_y + 1.19),
                         Inches(2.4), Inches(0.15),
                         text=ev["body"], size=9, color=muted, align="center",
                         line_spacing=1.3)
            add_line(slide, Inches(ex), Inches(timeline_y + dot_r),
                     Inches(ex), Inches(timeline_y + 0.40),
                     muted, weight=0.5)


def slide_equipe(prs, *, title, people, client="[NOME CLIENTE]",
                  date="MAIO 2026", theme="light"):
    """Cards de equipe.

    people: lista de {"name": str, "role": str, "skills": [str, str, str]}
    """
    slide = new_slide(prs, theme=theme)
    standard_frame(slide, eyebrow="EQUIPE · BIOS", title=title,
                    section="Equipe", client=client, date=date, theme=theme)

    fg = NIGHT_10 if theme == "light" else NIGHT_100
    container_bg = NIGHT_95 if theme == "light" else NIGHT_15
    muted = NIGHT_40 if theme == "light" else NIGHT_85

    n = len(people)
    card_w = (12.5 - (n - 1) * 0.25) / n
    card_h = 4.4
    start_y = 2.4

    for i, p in enumerate(people):
        x = 0.4 + i * (card_w + 0.25)
        y = start_y

        # Card bg
        add_rect(slide, Inches(x), Inches(y), Inches(card_w), Inches(card_h),
                 fill=container_bg, radius=0.04)

        # Avatar circle (placeholder)
        avatar_size = 1.6
        avatar_x = x + (card_w - avatar_size) / 2
        avatar_y = y + 0.35
        add_ellipse(slide, Inches(avatar_x), Inches(avatar_y),
                     Inches(avatar_size), Inches(avatar_size), fill=YELLOW)
        # Initial
        initials = "".join([w[0] for w in p["name"].split()[:2]])
        add_text(slide, Inches(avatar_x), Inches(avatar_y),
                 Inches(avatar_size), Inches(avatar_size),
                 text=initials, size=42, bold=True, color=NIGHT_10,
                 align="center", anchor="middle", is_head=True)

        # Nome
        add_text(slide, Inches(x + 0.2), Inches(y + 2.15),
                 Inches(card_w - 0.4), Inches(0.45),
                 text=p["name"], size=16, bold=True, color=fg,
                 align="center", is_head=True)
        # Role
        add_text(slide, Inches(x + 0.2), Inches(y + 2.65),
                 Inches(card_w - 0.4), Inches(0.35),
                 text=p["role"], size=11, color=muted,
                 align="center", letter_spacing=200, is_head=True)

        # Divider yellow
        div_y = y + 3.15
        add_line(slide, Inches(x + (card_w - 0.6) / 2), Inches(div_y),
                 Inches(x + (card_w + 0.6) / 2), Inches(div_y), YELLOW, weight=2)

        # Skills (até 3)
        for j, skill in enumerate(p.get("skills", [])[:3]):
            sy = y + 3.4 + j * 0.3
            add_text(slide, Inches(x + 0.3), Inches(sy),
                     Inches(card_w - 0.6), Inches(0.3),
                     text=f"· {skill}", size=10, color=muted,
                     align="center", line_spacing=1.3)


def slide_roadmap(prs, *, title, quarters, client="[NOME CLIENTE]",
                   date="MAIO 2026", theme="light"):
    """Roadmap visual Q1-Q4 (ou outros períodos).

    quarters: lista de {"label": "Q1 2026", "initiatives": [str, str, str], "highlight": bool}
    """
    slide = new_slide(prs, theme=theme)
    standard_frame(slide, eyebrow="ROADMAP · TRIMESTRAL", title=title,
                    section="Roadmap", client=client, date=date, theme=theme)

    fg = NIGHT_10 if theme == "light" else NIGHT_100
    container_bg = NIGHT_95 if theme == "light" else NIGHT_15
    muted = NIGHT_40 if theme == "light" else NIGHT_85

    n = len(quarters)
    card_w = (12.5 - (n - 1) * 0.2) / n
    card_h = 4.4
    start_y = 2.4

    for i, q in enumerate(quarters):
        x = 0.4 + i * (card_w + 0.2)
        y = start_y

        # 3 estados: "normal" | "highlight" (focal yellow) | "muted" (near-future slate)
        # Retrocompat: highlight=True ainda funciona
        state = q.get("state") or ("highlight" if q.get("highlight") else "normal")

        if state == "highlight":
            bg = YELLOW
            label_color = NIGHT_10
            sub_color = NIGHT_10
            bullet_color = NIGHT_95
            text_color = NIGHT_10
        elif state == "muted":
            bg = NIGHT_85 if theme == "light" else NIGHT_40
            label_color = NIGHT_10
            sub_color = NIGHT_40
            bullet_color = NIGHT_40
            text_color = NIGHT_10
        else:
            bg = container_bg
            label_color = YELLOW
            sub_color = muted
            bullet_color = YELLOW
            text_color = fg

        add_rect(slide, Inches(x), Inches(y), Inches(card_w), Inches(card_h),
                 fill=bg, radius=0.04)
        # NB: top stripe yellow removida (deltas v2 — limpar ruído)

        # Quarter label (height tight 0.22 — v2)
        add_text(slide, Inches(x + 0.25), Inches(y + 0.25),
                 Inches(card_w - 0.5), Inches(0.22),
                 text=q["label"].upper(), size=14, bold=True, color=label_color,
                 letter_spacing=300, is_head=True)

        if q.get("subtitle"):
            add_text(slide, Inches(x + 0.25), Inches(y + 0.65),
                     Inches(card_w - 0.5), Inches(0.3),
                     text=q["subtitle"], size=10, color=sub_color,
                     letter_spacing=100, is_head=True)

        add_line(slide, Inches(x + 0.25), Inches(y + 1.05),
                 Inches(x + card_w - 0.25), Inches(y + 1.05),
                 NIGHT_85 if theme == "light" else NIGHT_40, weight=1)

        for j, init in enumerate(q.get("initiatives", [])[:5]):
            iy = y + 1.25 + j * 0.6
            add_ellipse(slide, Inches(x + 0.3), Inches(iy + 0.13),
                         Inches(0.12), Inches(0.12), fill=bullet_color)
            add_text(slide, Inches(x + 0.55), Inches(iy),
                     Inches(card_w - 0.8), Inches(0.55),
                     text=init, size=11, color=text_color, line_spacing=1.35)


def slide_pricing(prs, *, title, plans, client="[NOME CLIENTE]",
                   date="MAIO 2026", theme="light"):
    """3 pricing tiers comparativos.

    plans: lista de {"name": str, "price": str, "subtitle": str (opcional),
                     "features": [str, ...], "highlight": bool, "cta": str}
    """
    slide = new_slide(prs, theme=theme)
    standard_frame(slide, eyebrow="OFERTA · PLANOS", title=title,
                    section="Investimento", client=client, date=date, theme=theme)

    fg = NIGHT_10 if theme == "light" else NIGHT_100
    container_bg = NIGHT_95 if theme == "light" else NIGHT_15
    muted = NIGHT_40 if theme == "light" else NIGHT_85

    n = len(plans)
    card_w = (12.5 - (n - 1) * 0.3) / n
    card_h = 4.5
    start_y = 2.2

    for i, plan in enumerate(plans):
        x = 0.4 + i * (card_w + 0.3)
        y = start_y
        is_highlight = plan.get("highlight", False)

        if is_highlight:
            # v2: highlight em YELLOW pleno (não mais YELLOW_SOFT — mais protagonista)
            bg = YELLOW
            offset_y = -0.15
        else:
            bg = container_bg
            offset_y = 0

        add_rect(slide, Inches(x), Inches(y + offset_y),
                 Inches(card_w), Inches(card_h - offset_y * 2),
                 fill=bg, radius=0.04)

        if is_highlight:
            # Badge "RECOMENDADO"
            badge_w = card_w - 0.4
            add_rect(slide, Inches(x + 0.2), Inches(y + offset_y - 0.2),
                     Inches(badge_w), Inches(0.32),
                     fill=NIGHT_10, radius=0.5)
            add_text(slide, Inches(x + 0.2), Inches(y + offset_y - 0.2),
                     Inches(badge_w), Inches(0.32),
                     text="MAIS ESCOLHIDO", size=9, bold=True, color=YELLOW,
                     align="center", anchor="middle", letter_spacing=400, is_head=True)

        # Plan name
        add_text(slide, Inches(x + 0.25), Inches(y + 0.3 + offset_y),
                 Inches(card_w - 0.5), Inches(0.4),
                 text=plan["name"].upper(), size=14, bold=True, color=muted,
                 letter_spacing=300, is_head=True)

        # Price (h=0.5 default — suficiente pra 40pt sem invadir subtitle)
        add_text(slide, Inches(x + 0.25), Inches(y + 0.75 + offset_y),
                 Inches(card_w - 0.5), Inches(0.5),
                 text=plan["price"], size=40, bold=True, color=fg,
                 is_head=True, line_spacing=0.95, letter_spacing=-15)

        # Subtitle
        if plan.get("subtitle"):
            add_text(slide, Inches(x + 0.25), Inches(y + 1.65 + offset_y),
                     Inches(card_w - 0.5), Inches(0.3),
                     text=plan["subtitle"], size=10, color=muted)

        # Divider
        add_line(slide, Inches(x + 0.25), Inches(y + 2.0 + offset_y),
                 Inches(x + card_w - 0.25), Inches(y + 2.0 + offset_y),
                 NIGHT_85 if theme == "light" else NIGHT_40, weight=1)

        # Features
        for j, feat in enumerate(plan.get("features", [])[:5]):
            fy = y + 2.2 + offset_y + j * 0.35
            add_text(slide, Inches(x + 0.25), Inches(fy),
                     Inches(0.3), Inches(0.3),
                     text="✓", size=12, bold=True, color=YELLOW, is_head=True)
            add_text(slide, Inches(x + 0.55), Inches(fy),
                     Inches(card_w - 0.8), Inches(0.3),
                     text=feat, size=10, color=fg, line_spacing=1.3)

        # CTA (no fim do card)
        cta_y = y + card_h + offset_y - 0.65
        cta_bg = NIGHT_10 if is_highlight else NIGHT_100
        cta_color = NIGHT_100 if is_highlight else NIGHT_10
        cta_line = None if is_highlight else NIGHT_10
        cta_pill = add_rect(slide, Inches(x + 0.25), Inches(cta_y),
                             Inches(card_w - 0.5), Inches(0.45),
                             fill=cta_bg, line=cta_line, radius=0.5)
        add_text(slide, Inches(x + 0.25), Inches(cta_y),
                 Inches(card_w - 0.5), Inches(0.45),
                 text=plan.get("cta", "Saiba mais").upper(), size=10, bold=True,
                 color=cta_color, align="center", anchor="middle",
                 letter_spacing=200, is_head=True)


def slide_funnel(prs, *, title, stages, client="[NOME CLIENTE]",
                   date="MAIO 2026", theme="light"):
    """Funil de conversão.

    stages: lista de {"name": str, "value": int, "conversion": str (opcional)}
    """
    slide = new_slide(prs, theme=theme)
    standard_frame(slide, eyebrow="CONVERSÃO · FUNIL", title=title,
                    section="Resultados", client=client, date=date, theme=theme)

    fg = NIGHT_10 if theme == "light" else NIGHT_100
    muted = NIGHT_40 if theme == "light" else NIGHT_85

    n = len(stages)
    max_value = max(s["value"] for s in stages)

    # Trapézio centralizado (funil afunila pra baixo)
    chart_x = 3.0
    chart_y = 2.3
    chart_w = 6.0
    chart_h = 4.3
    stage_h = chart_h / n

    for i, stage in enumerate(stages):
        # v2: decrescimento sublinear (expoente 0.4) — preserva legibilidade nos últimos
        # stages quando há queda forte (ex: 12k → 252 daria 0.02 linear, invisível)
        ratio = (stage["value"] / max_value) ** 0.4
        bar_w = chart_w * ratio
        bar_x = chart_x + (chart_w - bar_w) / 2
        bar_y = chart_y + i * stage_h

        # Cor: primeiro stage YELLOW, depois gradiente pra muted
        if i == 0:
            color = YELLOW
        elif i == n - 1:
            color = YELLOW_SOFT
        else:
            color = NIGHT_85 if theme == "light" else NIGHT_40

        add_rect(slide, Inches(bar_x), Inches(bar_y),
                 Inches(bar_w), Inches(stage_h - 0.1),
                 fill=color, radius=0.02)

        # Label dentro do stage
        add_text(slide, Inches(bar_x), Inches(bar_y),
                 Inches(bar_w), Inches(stage_h - 0.1),
                 text=stage["name"], size=12, bold=True, color=NIGHT_10,
                 align="center", anchor="middle", is_head=True)

        # Number à direita
        add_text(slide, Inches(chart_x + chart_w + 0.3), Inches(bar_y + 0.05),
                 Inches(2.0), Inches(0.5),
                 text=f"{stage['value']:,}".replace(",", "."),
                 size=22, bold=True, color=fg,
                 is_head=True, letter_spacing=-10)

        # Conversion rate label (% de queda)
        if stage.get("conversion") and i > 0:
            add_text(slide, Inches(chart_x + chart_w + 0.3), Inches(bar_y + 0.5),
                     Inches(2.0), Inches(0.3),
                     text=stage["conversion"], size=10, color=muted,
                     letter_spacing=100, is_head=True)

        # Conector vertical tracejado entre stages
        if i < n - 1:
            # Já natural pelo espaçamento stage_h
            pass


def slide_faq(prs, *, title, questions, client="[NOME CLIENTE]",
                date="MAIO 2026", theme="light"):
    """FAQ — perguntas + respostas.

    questions: lista de {"q": str, "a": str}
    """
    slide = new_slide(prs, theme=theme)
    standard_frame(slide, eyebrow="DÚVIDAS · FAQ", title=title,
                    section="FAQ", client=client, date=date, theme=theme)

    fg = NIGHT_10 if theme == "light" else NIGHT_100
    muted = NIGHT_40 if theme == "light" else NIGHT_85

    n = min(len(questions), 4)
    item_h = 4.5 / n
    start_y = 2.3

    for i, qa in enumerate(questions[:n]):
        y = start_y + i * item_h

        # Number badge yellow
        add_text(slide, Inches(0.4), Inches(y),
                 Inches(0.8), Inches(0.4),
                 text=f"0{i+1}", size=24, bold=True, color=YELLOW,
                 is_head=True, letter_spacing=-15)

        # Question (bold)
        add_text(slide, Inches(1.3), Inches(y),
                 Inches(11.5), Inches(0.45),
                 text=qa["q"], size=15, bold=True, color=fg,
                 is_head=True, line_spacing=1.2)

        # Answer (regular) — v2: gap tight (0.37 ao invés de 0.55) — adesão à pergunta
        add_text(slide, Inches(1.3), Inches(y + 0.37),
                 Inches(11.5), Inches(item_h - 0.52),
                 text=qa["a"], size=11, color=muted, line_spacing=1.5)

        # Divider entre items
        if i < n - 1:
            div_y = y + item_h - 0.1
            add_line(slide, Inches(1.3), Inches(div_y),
                     Inches(12.9), Inches(div_y),
                     NIGHT_85 if theme == "light" else NIGHT_40, weight=0.5)


# ====================================================================
# BUILDERS v3 — FIEL ao PPT Modelo Geral
# Fonte: design/metta-pptx-canonico.md (análise shape-a-shape)
# ====================================================================

def slide_capa(prs, *, statement, subline, client="[NOME CLIENTE]",
                date="Maio 2026", theme="dark", date_label="Maio 2026"):
    """CAPA institucional (modelo geral slide 1/2).

    statement: title big (size 55pt). Aceita markdown `**bold**` pra ênfase editorial
               (ex: `Somos especialistas **em resultados.**`)
    subline:   subtítulo Inter 10pt embaixo do statement
    theme:     'dark' (BG_DARK) | 'light' (BG_LIGHT_COOL)

    Layout do modelo: title bloco lado esquerdo (x=0.54 a 7.32) + picture
    ornamental ocupa direita (placeholder se sem foto).
    NÃO usa anéis concêntricos (anti-padrão §6.4).
    """
    bg_theme = "dark" if theme == "dark" else "light-cool"
    slide = new_slide(prs, theme=bg_theme)
    fg_bold = TEXT_DARK_BOLD if theme == "dark" else TEXT_LIGHT_BOLD
    muted = TEXT_DARK_MUTED if theme == "dark" else TEXT_LIGHT_MUTED
    subline_color = TEXT_DARK_MUTED if theme == "dark" else TEXT_LIGHT_SOFT
    arrow_color = TEXT_DARK_PURE if theme == "dark" else muted

    # Header COMPACT: logo mini + date micro top-right
    add_header(slide, theme=theme, mode="compact")

    # Ornamento lateral direito — placeholder (sub. depois por picture real)
    # Bloco yellow geométrico discreto (não anel concêntrico)
    add_rect(slide, Inches(8.5), Inches(0.0), Inches(4.833), Inches(7.5),
             fill=YELLOW if theme == "light" else CONTAINER_DARK, radius=0)

    # HEADLINE big (lado esquerdo)
    add_title_split(slide, x=Inches(0.54), y=Inches(2.616),
                     w=Inches(6.784), text=statement,
                     size=55, color=fg_bold)

    # SUBLINE (Inter 10pt)
    add_text(slide, Inches(0.67), Inches(4.464), Inches(5.21), Inches(0.417),
             text=subline, size=10, color=subline_color, line_spacing=1.5)

    # Footer COMPACT
    add_footer(slide, date=date_label.upper(), client=client, section="INSTITUCIONAL",
               theme=theme, mode="compact")
    return slide


def slide_section_divider(prs, *, title, progress=None, theme="dark"):
    """SECTION DIVIDER (modelo geral slide 10).

    title:    título grande centralizado (size 72pt). Sem yellow line accent.
    progress: opcional "01 / 04" no canto esquerdo
    theme:    'dark' (default, mais comum no modelo) | 'light' | 'yellow' (slide 3)

    Sem header eyebrow (clean transition). Footer COMPACT.
    """
    slide = new_slide(prs, theme=theme)
    fg = TEXT_DARK_PURE if theme == "dark" else TEXT_LIGHT_BOLD
    muted = TEXT_DARK_MUTED if theme == "dark" else TEXT_LIGHT_MUTED

    add_logo(slide, theme=theme, mode="main")

    if progress:
        add_text(slide, Inches(0.4), Inches(0.72), Inches(3), Inches(0.25),
                 text=progress, size=9, bold=True, color=muted,
                 is_head=True)

    # Title centralizado y=3.169 h=1.151 size 72pt (modelo exato)
    add_title_split(slide, x=Inches(0.4), y=Inches(3.169),
                     w=Inches(12.5), text=title.title(),
                     size=72, color=fg)
    # Centraliza horizontalmente (default LEFT — re-align via XML)
    # Truque: add_title_split usa LEFT por default. Pra centralizar, sobrescrevemos
    # o paragraph alignment — mas como add_title_split é interno, é mais simples
    # usar add_text como title centralizado:
    # Mas perderia o split bold/regular. Aceitamos LEFT por enquanto.

    add_footer(slide, theme=theme, section=title.upper(), mode="compact",
               client="[NOME CLIENTE]")
    return slide


def slide_quem_somos_v3(prs, *, title="**Quem** somos", body, stats=None,
                          client="[NOME CLIENTE]", date="MAIO 2026"):
    """QUEM SOMOS split column (modelo geral slide 4).

    Layout DUAS COLUNAS:
    - Esquerda (x=0 a 4.785): texto sobre dark
    - Direita (x=4.986 a 13.333): bloco yellow + foto editorial bleed

    title: aceita markdown `**bold**` para palavras-chave (default: `**Quem** somos`)
    """
    slide = new_slide(prs, theme="dark")
    muted = TEXT_DARK_MUTED

    # Header SPLIT (logo + eyebrow só esquerda + divider parcial)
    add_header(slide, eyebrow="PROPOSTA · 2026", client=client,
               theme="dark", mode="split")

    # COLUNA DIREITA: bloco yellow geometric (background)
    add_rect(slide, Inches(4.986), Inches(0.0), Inches(10.152), Inches(7.799),
             fill=YELLOW, radius=0)
    # Placeholder pra foto editorial sobre o yellow
    add_text(slide, Inches(7.5), Inches(3.5), Inches(4), Inches(0.5),
             text="[ FOTO EDITORIAL ]", size=10, color=TEXT_LIGHT_BOLD,
             align="center", anchor="middle", is_head=True)

    # COLUNA ESQUERDA: title grande + subline
    add_title_split(slide, x=Inches(0.556), y=Inches(2.525),
                     w=Inches(4.034), text=title, size=70, color=TEXT_DARK_PURE)

    add_text(slide, Inches(0.556), Inches(4.604), Inches(3.771), Inches(0.812),
             text=body, size=13, color=muted, line_spacing=1.5)

    # Stats (opcional) — 1 coluna esquerda, abaixo do subline
    if stats:
        for i, (num, label) in enumerate(stats[:3]):
            sy = 5.5 + i * 0.55
            add_text(slide, Inches(0.556), Inches(sy), Inches(2.5), Inches(0.3),
                     text=num, size=16, bold=True, color=YELLOW,
                     is_head=True)
            add_text(slide, Inches(2.3), Inches(sy + 0.05), Inches(2.5), Inches(0.25),
                     text=label, size=9, color=muted, is_head=True)

    # Footer SPLIT
    add_footer(slide, date=date, client=client, section="CLIENTE",
               theme="dark", mode="split")
    return slide


def slide_prova_big_number_v3(prs, *, eyebrow_main, big_number, big_unit, body=None,
                                 caption=None, client="[NOME CLIENTE]",
                                 date="MAIO 2026"):
    """PROVA SOCIAL big number (modelo geral slide 5).

    Padrão fiel:
    - BG dark, header MAIN + footer MAIN
    - Eyebrow body Zalando SemiExpanded 24.7pt cor #B0CAD8
    - Big number: 165pt bold yellow (valor) + 175pt regular branco (unidade) — sizes DIFERENTES
    - Picture bleed à direita (placeholder)

    eyebrow_main: ex "Em metas batidas nos\\núltimos 12 meses:"
    big_number:   ex "R$ 8,2"
    big_unit:     ex "bi"
    """
    slide = new_slide(prs, theme="dark")
    muted = TEXT_DARK_MUTED

    add_header(slide, eyebrow="PROVA DE MERCADO · 12 MESES", client=client,
               theme="dark", mode="main")

    # Picture placeholder direita (bleed)
    add_rect(slide, Inches(7.167), Inches(0.583), Inches(6.264), Inches(7.0),
             fill=CONTAINER_DARK, radius=0.04)
    add_text(slide, Inches(7.167), Inches(3.5), Inches(6.264), Inches(0.5),
             text="[ IMAGEM / GRÁFICO ]", size=10, color=muted,
             align="center", anchor="middle", is_head=True)

    # Eyebrow body (SemiExpanded)
    add_text(slide, Inches(0.4), Inches(3.276), Inches(5.565), Inches(0.779),
             text=eyebrow_main, size=24, bold=True, color=muted,
             font_name=FONT_HEADS_SEMI, line_spacing=1.2, is_head=True)

    # BIG NUMBER (valor — bold yellow 165pt)
    add_text(slide, Inches(0.4), Inches(4.472), Inches(9.62), Inches(2.351),
             text=big_number, size=165, bold=True, color=YELLOW,
             line_spacing=0.9, letter_spacing=-40, is_head=True)

    # BIG NUMBER (unidade — REGULAR branco 175pt)
    add_text(slide, Inches(10.135), Inches(4.472), Inches(2.372), Inches(2.351),
             text=big_unit, size=175, bold=False, color=TEXT_DARK_PURE,
             line_spacing=0.9, letter_spacing=-40, is_head=True)

    # Caption secundário (opcional)
    if caption:
        add_text(slide, Inches(9.655), Inches(3.352), Inches(3.21), Inches(0.534),
                 text=caption, size=12, color=muted, line_spacing=1.4)

    add_footer(slide, date=date, client=client, section="INSTITUCIONAL",
               theme="dark", mode="main")
    return slide


def slide_encerramento_v3(prs, *, title, body, info_cards, cta,
                            client="[NOME CLIENTE]", date="MAIO 2026"):
    """ENCERRAMENTO CTA (modelo geral slide 19).

    Sem anel decorativo bottom-right (anti-padrão §6.4).
    """
    slide = new_slide(prs, theme="light")
    fg = TEXT_LIGHT_BOLD
    muted = TEXT_LIGHT_MUTED

    add_header(slide, eyebrow="PRÓXIMO PASSO", client=client,
               theme="light", mode="main")

    # Title central
    add_title_split(slide, x=Inches(0.4), y=Inches(1.5),
                     w=Inches(12.5), text=title, size=56, color=fg)

    # Body
    add_text(slide, Inches(2.0), Inches(3.2), Inches(9.333), Inches(0.8),
             text=body, size=15, color=muted, align="center", line_spacing=1.5)

    # Info cards (sem anel decorativo)
    n = len(info_cards)
    card_w = (10 - (n - 1) * 0.2) / n
    for i, card in enumerate(info_cards):
        x = 1.667 + i * (card_w + 0.2)
        add_rect(slide, Inches(x), Inches(4.3), Inches(card_w), Inches(1.3),
                 fill=CONTAINER_LIGHT, radius=0.04)
        add_text(slide, Inches(x), Inches(4.5), Inches(card_w), Inches(0.3),
                 text=card["label"].upper(), size=10, bold=True,
                 color=YELLOW, align="center", is_head=True)
        add_text(slide, Inches(x), Inches(4.9), Inches(card_w), Inches(0.5),
                 text=card["value"], size=14, bold=True, color=fg,
                 align="center", is_head=True)

    # CTA pill
    cta_w = 3.4
    cta_x = (13.333 - cta_w) / 2
    add_rect(slide, Inches(cta_x), Inches(6.0), Inches(cta_w), Inches(0.65),
             fill=TEXT_LIGHT_BOLD, radius=0.5)
    add_text(slide, Inches(cta_x), Inches(6.0), Inches(cta_w), Inches(0.65),
             text=cta.upper(), size=14, bold=True, color=TEXT_DARK_PURE,
             align="center", anchor="middle", is_head=True)

    add_footer(slide, date=date, client=client, section="FECHAMENTO",
               theme="light", mode="main")
    return slide
